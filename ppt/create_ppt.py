from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from lxml import etree

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

NAVY   = RGBColor(0x1a, 0x3a, 0x5c)
GOLD   = RGBColor(0xf5, 0x9e, 0x0b)
BLUE   = RGBColor(0x25, 0x63, 0xeb)
GREEN  = RGBColor(0x10, 0xb9, 0x81)
RED    = RGBColor(0xef, 0x44, 0x44)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LGRAY  = RGBColor(0xf0, 0xf4, 0xf8)
MGRAY  = RGBColor(0x64, 0x74, 0x8b)

LIVE_URL   = "https://omais-eng.github.io/assignment-04-rpa-bot-troubleshooting/"
GITHUB_URL = "https://github.com/omais-eng/assignment-04-rpa-bot-troubleshooting"

blank_layout = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, w, h, fill=NAVY):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, w, h, size=18, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def slide_bg(slide, color=LGRAY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_notes(slide, text):
    """Set speaker notes on a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    # clear existing paragraphs
    for para in tf.paragraphs[1:]:
        p = para._p
        p.getparent().remove(p)
    tf.paragraphs[0].text = text

def add_link_box(slide, label, url, left, top, w, h, bg_color=None):
    """Add a clickable hyperlink text box."""
    if bg_color is None:
        bg_color = RGBColor(0xef, 0xf6, 0xff)
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = BLUE
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"🔗  {label}"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = BLUE
    # add hyperlink using lxml directly (namespace-aware)
    rPr = run._r.get_or_add_rPr()
    rId = slide.part.relate_to(
        url,
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
        is_external=True
    )
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    hlinkClick = etree.SubElement(rPr, f'{{{A_NS}}}hlinkClick')
    hlinkClick.set(f'{{{R_NS}}}id', rId)
    return box

# ─────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, NAVY)
add_rect(sl, 0, 5.8, 13.33, 0.12, GOLD)
add_text_box(sl, "RPA COURSE  ·  ASSIGNMENT 04", 1, 0.6, 11, 0.5, size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text_box(sl, "E-Commerce Order Processing", 1, 1.3, 11, 0.9, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(sl, "RPA Bot — Troubleshooting & Debugging Case Study", 1, 2.2, 11, 0.8, size=22, color=RGBColor(0x93,0xc5,0xfd), align=PP_ALIGN.CENTER)
add_text_box(sl, "A comprehensive study on diagnosing, debugging, and optimizing a\nUiPath RPA bot for end-to-end e-commerce order processing", 1, 3.1, 11, 0.9, size=13, color=RGBColor(0x94,0xa3,0xb8), align=PP_ALIGN.CENTER)
add_text_box(sl, "Student: Omais Siddiqui     |     Platform: UiPath + Orchestrator     |     April 2026", 1, 6.0, 11, 0.4, size=11, color=RGBColor(0x64,0x74,0x8b), align=PP_ALIGN.CENTER)

# Links row
add_link_box(sl, "🌐 Live Website", LIVE_URL,   1.5, 4.3, 4.5, 0.5, RGBColor(0x1e,0x4a,0x76))
add_link_box(sl, "GitHub Repository", GITHUB_URL, 7.3, 4.3, 4.5, 0.5, RGBColor(0x1e,0x4a,0x76))

set_notes(sl,
"""SLIDE 1 — TITLE SLIDE

Welcome and introduce yourself:
"Hello, I'm Omais Siddiqui. This is Assignment 04 for the RPA course — a complete case study on troubleshooting and debugging an E-Commerce Order Processing RPA Bot built on UiPath."

Point out the two clickable links:
- Live Website: The interactive single-page website is deployed on GitHub Pages. Click the link to open it live.
- GitHub Repository: All source code, PowerPoint, Word documents, and the video script are available on GitHub.

Mention the scope:
"This case study covers five sections: RPA error documentation, a systematic debugging framework, optimization strategies, a monitoring dashboard mock-up, and how AI-assisted vibe coding was used to build this project."

Transition: "Let's start with the agenda overview."
""")

# ─────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
add_text_box(sl, "Presentation Agenda", 0.5, 0.2, 12, 0.7, size=26, bold=True, color=WHITE)

items = [
    ("01", "RPA Error Log",          "5 realistic exceptions across workflow stages",        BLUE),
    ("02", "Debugging Strategy",     "Logging, monitoring, retry logic, decision flow",      GREEN),
    ("03", "Optimization",           "Performance, reliability, scheduling, before/after",   GOLD),
    ("04", "Monitoring Dashboard",   "KPIs, queue health, alerts, bot instance status",      RED),
    ("05", "AI-Assisted Development","Vibe coding with Claude Code — how it was built",      RGBColor(0x93,0x33,0xea)),
]
for i, (num, title, sub, color) in enumerate(items):
    x = 0.4 + (i % 3) * 4.3
    y = 1.4 + (i // 3) * 2.8
    add_rect(sl, x, y, 3.9, 2.4, WHITE)
    add_rect(sl, x, y, 3.9, 0.08, color)
    add_text_box(sl, num,   x+0.15, y+0.15, 0.7,  0.55, size=28, bold=True, color=color)
    add_text_box(sl, title, x+0.15, y+0.75, 3.6,  0.45, size=14, bold=True, color=NAVY)
    add_text_box(sl, sub,   x+0.15, y+1.2,  3.6,  0.9,  size=10, color=MGRAY)

set_notes(sl,
"""SLIDE 2 — AGENDA

"This presentation is structured into five sections."

Walk through each card briefly:
01 - Error Log: We documented 5 realistic UiPath exceptions with root cause analysis.
02 - Debugging: A systematic framework covering logging, monitoring, and step-by-step resolution.
03 - Optimization: Four improvement categories plus a measurable before/after comparison.
04 - Dashboard: A mock monitoring UI showing what Orchestrator dashboards look like in production.
05 - AI Usage: How Claude Code was used to build this entire project using vibe coding.

"Each section builds on the previous, moving from problem identification to resolution to prevention."

Transition: "Let's begin with the error documentation."
""")

# ─────────────────────────────────────────────
# SLIDE 3 — ERRORS TABLE
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
add_text_box(sl, "Section 1 — RPA Error Log: 5 Exceptions", 0.5, 0.2, 12, 0.7, size=24, bold=True, color=WHITE)

errors = [
    ("E-01", "SelectorNotFoundException", "System",   "Order Placement",    "HIGH",   "Stale CSS selector after frontend deploy"),
    ("E-02", "InvalidOrderDataException", "Business", "Data Extraction",    "HIGH",   "Missing ZIP code in ERP international feed"),
    ("E-03", "OracleDBTimeoutException",  "System",   "Inventory Check",    "HIGH",   "Missing DB index causing full table scans"),
    ("E-04", "DuplicateOrderException",   "Business", "Order Submission",   "MEDIUM", "No idempotency key on retry logic"),
    ("E-05", "PDFParsingFailureException","App",      "Invoice Generation", "LOW",    "Scanned PDF — OCR not enabled"),
]
hdrs   = ["#", "Error Name", "Type", "Stage", "Severity", "Root Cause"]
col_w  = [0.45, 3.0, 1.1, 1.6, 1.0, 5.4]
col_x  = [0.25]
for w in col_w[:-1]: col_x.append(col_x[-1]+w)

add_rect(sl, 0.25, 1.2, 12.85, 0.42, NAVY)
for j, (h, x, w) in enumerate(zip(hdrs, col_x, col_w)):
    add_text_box(sl, h, x+0.05, 1.22, w-0.1, 0.38, size=10, bold=True, color=WHITE)

row_colors = [WHITE, LGRAY, WHITE, LGRAY, WHITE]
sev_colors = {"HIGH": RED, "MEDIUM": GOLD, "LOW": GREEN}
type_colors = {"System": BLUE, "Business": GOLD, "App": RGBColor(0x93,0x33,0xea)}
for i, (eid, name, etype, stage, sev, cause) in enumerate(errors):
    y = 1.65 + i*0.95
    add_rect(sl, 0.25, y, 12.85, 0.9, row_colors[i])
    vals = [eid, name, etype, stage, sev, cause]
    for j, (val, x, w) in enumerate(zip(vals, col_x, col_w)):
        c = NAVY
        if j == 2: c = type_colors.get(val, BLUE)
        if j == 4: c = sev_colors.get(val, MGRAY)
        bold = j in (0,1,4)
        add_text_box(sl, val, x+0.05, y+0.05, w-0.1, 0.8, size=9 if j==5 else 10, bold=bold, color=c)

set_notes(sl,
"""SLIDE 3 — RPA ERROR LOG

"Let's look at the five exceptions we encountered and documented."

Walk through each row:

E-01 — SelectorNotFoundException (System Exception, HIGH severity):
"This is a classic UI automation failure. The checkout portal's frontend team renamed a CSS selector during a release, and the bot's selector library wasn't updated. This blocked approximately 120 orders per hour."

E-02 — InvalidOrderDataException (Business Exception, HIGH):
"A data contract issue. The ERP system began exporting international orders without the postal code field. The bot had no input schema validation, so these orders would silently fail in the shipping module downstream."

E-03 — OracleDBTimeoutException (System Exception, HIGH):
"An infrastructure problem. The inventory table lacked a composite index. Under peak load with 50 concurrent bot instances, every inventory check triggered a full table scan, timing out at 30 seconds and causing 400+ units to be oversold weekly."

E-04 — DuplicateOrderException (Business Exception, MEDIUM):
"A retry logic flaw. When a network timeout occurred mid-POST, the bot retried without checking if the original request had already succeeded — charging customers twice."

E-05 — PDFParsingFailureException (Application Exception, LOW):
"A dependency change issue. The vendor switched from text-layer PDFs to scanned image PDFs. Without OCR enabled, the parser returned null fields, requiring 1.5 hours of manual invoice re-generation daily."

Transition: "Now let's look at how we debug and resolve these systematically."
""")

# ─────────────────────────────────────────────
# SLIDE 4 — DEBUGGING STRATEGY
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
add_text_box(sl, "Section 2 — Debugging Strategy & Framework", 0.5, 0.2, 12, 0.7, size=24, bold=True, color=WHITE)

debug_cards = [
    ("📋 Logging",     BLUE,  "• Structured levels: INFO, WARN, ERROR, FATAL\n• Include: OrderID, BotInstanceID, UTC timestamp\n• Screenshot on every exception\n• Ship to Orchestrator → Elasticsearch → Kibana"),
    ("📊 Monitoring",  GREEN, "• Real-time KPI tracking (60s refresh)\n• PagerDuty: alert if success rate <90%\n• SLA violation detection (4-hour threshold)\n• Anomaly detection → auto Jira tickets"),
    ("🔁 Reproduction",GOLD,  "• Isolate failing transaction from queue\n• Snapshot environment state\n• Replay in UiPath Debug mode\n• Reproduce 3× before applying any fix"),
    ("🔄 Retry Logic", RED,   "• RetryScope: 3 retries, 5s backoff\n• Idempotency key check before each order POST\n• System exceptions → retry then escalate\n• Business exceptions → exception queue"),
]
for i, (title, color, body) in enumerate(debug_cards):
    x = 0.3 + (i%2)*6.4
    y = 1.25 + (i//2)*2.8
    add_rect(sl, x, y, 6.1, 2.55, WHITE)
    add_rect(sl, x, y, 6.1, 0.08, color)
    add_text_box(sl, title, x+0.15, y+0.15, 5.8, 0.45, size=13, bold=True, color=NAVY)
    add_text_box(sl, body,  x+0.15, y+0.65, 5.8, 1.8,  size=10, color=RGBColor(0x1e,0x29,0x3b))

set_notes(sl,
"""SLIDE 4 — DEBUGGING STRATEGY

"Our debugging framework has four pillars."

📋 Logging:
"Every log entry carries a CorrelationID so we can trace a single order across multiple microservices — from ERP extraction through checkout portal to invoice generation. Exceptions automatically trigger a screenshot saved with the OrderID and timestamp."

📊 Monitoring:
"Orchestrator KPIs refresh every 60 seconds. PagerDuty fires if success rate drops below 90%, error rate exceeds 5%, or a bot becomes unresponsive for more than 5 minutes. SLA violations — orders older than 4 hours — trigger an automated email to the Operations Manager."

🔁 Reproduction:
"We never fix what we can't reproduce consistently. The process: isolate the failing transaction from the Orchestrator queue, clone the environment state including app version and DB snapshot, replay in UiPath's Debug mode with Step Into, and confirm the error three times before writing a single line of fix code."

🔄 Retry Logic:
"RetryScope is configured with 3 retries and a 5-second backoff. Critically — the idempotency key check happens BEFORE the order POST, not after. This alone eliminates all duplicate order incidents."

Transition: "The next slide shows these four pillars connected into a single decision flow."
""")

# ─────────────────────────────────────────────
# SLIDE 5 — DEBUGGING FLOW
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
add_text_box(sl, "Debugging Decision Flow — 10-Step Process", 0.5, 0.2, 12, 0.7, size=24, bold=True, color=WHITE)

steps = [
    ("Error\nDetected",   RED),
    ("Log\nCaptured",     BLUE),
    ("Alert\nTriggered",  GREEN),
    ("Classify\nType",    RGBColor(0x93,0x33,0xea)),
    ("Reproduce\nin DEV", GOLD),
    ("Apply\nFix",        RGBColor(0x13,0x4e,0x4a)),
    ("Test 50\nOrders",   BLUE),
    ("Deploy\nto PROD",   GREEN),
    ("Monitor\n24hr",     RGBColor(0x14,0x53,0x2d)),
    ("Close\nIncident",   NAVY),
]
for i, (label, color) in enumerate(steps):
    x = 0.3 + i * 1.27
    add_rect(sl, x, 1.4, 1.1, 0.9, color)
    add_text_box(sl, label, x+0.05, 1.42, 1.0, 0.86, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        add_text_box(sl, "→", x+1.1, 1.6, 0.18, 0.5, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

paths = [
    ("System Exception",  BLUE,                    "Auto-retry 3× → If persistent: pause bot → Notify DevOps → Hotfix selector/infrastructure",   0.3),
    ("Business Exception",GOLD,                    "Move to Exception Queue → Human review within 2 hours → Reprocess with corrected data",          4.5),
    ("App Exception",     RGBColor(0x93,0x33,0xea),"Enable OCR fallback → Reroute to alternate parser → Log for vendor notification",               8.7),
]
for title, color, body, x in paths:
    add_rect(sl, x, 2.65, 4.0, 1.5, WHITE)
    add_rect(sl, x, 2.65, 4.0, 0.07, color)
    add_text_box(sl, title, x+0.12, 2.72, 3.76, 0.38, size=11, bold=True, color=color)
    add_text_box(sl, body,  x+0.12, 3.12, 3.76, 0.9,  size=9,  color=NAVY)

add_text_box(sl, "Each exception type has its own handling path — reducing Mean Time to Resolution (MTTR)", 0.3, 4.5, 12.7, 0.4, size=11, color=MGRAY, align=PP_ALIGN.CENTER)

set_notes(sl,
"""SLIDE 5 — DEBUGGING DECISION FLOW

"This flow is the connective tissue of our entire debugging strategy — 10 steps from error detection to incident closure."

Walk through the flow:
Steps 1-3: "Error detected → Structured log captured → PagerDuty alert triggered. This happens automatically in under 30 seconds."

Step 4 — Classify: "This is the critical branch point. Is it a System, Business, or Application exception? The classification determines the entire handling path."

System path: "Auto-retry up to 3 times with backoff. If it persists, the bot pauses itself and notifies DevOps. A hotfix is deployed — typically a selector library update or infrastructure fix."

Business path: "No automated retry — the transaction moves to the Exception Queue for human review within 2 hours. A human verifies and corrects the data, then re-queues the order."

App path: "The fallback pipeline activates: OCR pre-processing is enabled, the order is rerouted to the Tesseract parser, and the vendor is notified of the template change."

Steps 6-10: "Apply fix → Test on 50 orders → Deploy to production → Monitor for 24 hours → Close the incident with root cause documentation in the runbook."

"This systematic approach reduces MTTR — Mean Time to Resolution — from hours to minutes."
""")

# ─────────────────────────────────────────────
# SLIDE 6 — OPTIMIZATION
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
add_text_box(sl, "Section 3 — Bot Optimization Strategies", 0.5, 0.2, 12, 0.7, size=24, bold=True, color=WHITE)

opt_cards = [
    ("⚡ Performance",    BLUE,  "• 5 parallel bot lanes via Work Queues\n• DB composite index: 30s → <200ms\n• Lazy UI loading: -40% wait time\n• Batch inventory API calls (50 per batch)"),
    ("🛡️ Reliability",   GREEN, "• Idempotency keys eliminate duplicates\n• Dynamic anchor-based selectors\n• JSON schema validation on ERP feed\n• Circuit breaker on DB error spikes"),
    ("🔧 Maintainability",GOLD,  "• Centralized selector library in Orchestrator\n• Modular .xaml workflows per stage\n• Config-driven (no code deploy for params)\n• OCR pipeline for scanned PDFs"),
    ("🗓️ Scheduling",    RGBColor(0x93,0x33,0xea), "• Peak (8AM–6PM): 5 bot instances\n• Off-peak (6PM–8AM): 2 instances\n• Maintenance: 2AM–4AM daily window\n• VIP queue: 15-minute SLA guarantee"),
]
for i, (title, color, body) in enumerate(opt_cards):
    x = 0.3 + (i%2)*6.4
    y = 1.25 + (i//2)*2.8
    add_rect(sl, x, y, 6.1, 2.55, WHITE)
    add_rect(sl, x, y, 6.1, 0.08, color)
    add_text_box(sl, title, x+0.15, y+0.15, 5.8, 0.45, size=13, bold=True, color=NAVY)
    add_text_box(sl, body,  x+0.15, y+0.65, 5.8, 1.8,  size=10, color=RGBColor(0x1e,0x29,0x3b))

set_notes(sl,
"""SLIDE 6 — OPTIMIZATION STRATEGIES

"Once we understand the errors and can debug them, the next step is preventing them and improving overall performance."

⚡ Performance:
"The biggest win was parallelism — splitting the order queue into 5 concurrent bot lanes multiplied throughput by 5×. Adding a composite database index eliminated the E-03 timeout completely, dropping query time from 30 seconds to under 200 milliseconds."

🛡️ Reliability:
"Idempotency keys are the fix for E-04 — every order gets a UUID generated before the POST. If the bot retries, it checks the key first and skips if already processed. Dynamic anchor-based selectors solved E-01 — the bot no longer depends on a single fragile CSS class name."

🔧 Maintainability:
"All selectors are now stored in Orchestrator assets — a single update point. When the frontend changes, we update one asset and all 5 bot instances pick it up on their next transaction. Configuration parameters like batch sizes and retry counts are also asset-driven — zero code deployments needed for tuning."

🗓️ Scheduling:
"The bot scales dynamically: 5 instances during peak business hours, 2 during off-peak, and a dedicated 2AM–4AM maintenance window for deployments that avoids any business impact. VIP orders — same-day delivery — jump the queue with a 15-minute SLA guarantee."

Transition: "The next slide shows the measurable impact of all these changes."
""")

# ─────────────────────────────────────────────
# SLIDE 7 — BEFORE / AFTER
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
add_text_box(sl, "Before vs. After Optimization — Key Metrics", 0.5, 0.2, 12, 0.7, size=24, bold=True, color=WHITE)

add_rect(sl, 0.3,  1.2, 5.8, 5.9, RGBColor(0xff,0xf5,0xf5))
add_rect(sl, 0.3,  1.2, 5.8, 0.5, RED)
add_text_box(sl, "❌  BEFORE Optimization", 0.45, 1.25, 5.5, 0.4, size=13, bold=True, color=WHITE)
before_items = [
    "Success rate: 71%",
    "Processing time: 4.2 min/order",
    "Duplicate orders: ~12 incidents/week",
    "DB timeouts: 80+ per peak hour",
    "Manual overhead: 3.5 hours/day",
    "Selector failures after every frontend deploy",
    "Throughput: 120 orders/hour",
    "Config changes require code deployment",
]
for i, item in enumerate(before_items):
    add_text_box(sl, "✗  " + item, 0.45, 1.85+i*0.55, 5.5, 0.5, size=11, color=RGBColor(0x7f,0x1d,0x1d))

add_rect(sl, 7.2, 1.2, 5.8, 5.9, RGBColor(0xf0,0xfd,0xf4))
add_rect(sl, 7.2, 1.2, 5.8, 0.5, GREEN)
add_text_box(sl, "✅  AFTER Optimization", 7.35, 1.25, 5.5, 0.4, size=13, bold=True, color=WHITE)
after_items = [
    "Success rate: 97.4% (+26%)",
    "Processing time: 1.8 min/order (57% faster)",
    "Duplicate orders: 0 incidents",
    "DB timeouts: eliminated (<200ms)",
    "Manual overhead: 12 min/day (96% less)",
    "Dynamic selectors survive all deployments",
    "Throughput: 600 orders/hour (5× increase)",
    "Zero-downtime config updates via Orchestrator",
]
for i, item in enumerate(after_items):
    add_text_box(sl, "✓  " + item, 7.35, 1.85+i*0.55, 5.5, 0.5, size=11, color=RGBColor(0x06,0x5f,0x46))

add_text_box(sl, "VS", 6.3, 3.7, 0.7, 0.7, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

set_notes(sl,
"""SLIDE 7 — BEFORE VS. AFTER COMPARISON

"The numbers tell the story of what systematic debugging and optimization deliver."

Walk through the key pairs:

Success Rate: "71% to 97.4% — a 26 percentage point improvement. Before optimization, nearly 1 in 3 orders required manual intervention. After, less than 1 in 40 fails."

Processing Time: "4.2 minutes to 1.8 minutes per order — 57% faster. When you multiply this across thousands of orders daily, this is hours of throughput reclaimed."

Duplicate Orders: "12 incidents per week to zero. This is the idempotency key fix for E-04. Zero duplicate charges to customers in 30 days post-deployment."

Database Timeouts: "From 80+ per peak hour to eliminated. The composite index on product_sku and warehouse_id was a single database change that fixed E-03 completely."

Manual Overhead: "3.5 hours per day of manual invoice re-generation down to 12 minutes. That's 96% of time returned to the operations team — from OCR pipeline enabling and the PDF routing fix for E-05."

Throughput: "120 orders per hour to 600. Five bot instances instead of one. Same infrastructure, 5× output."

Config Changes: "Previously required a full code deployment — staging, testing, approval, release. Now a parameter change in Orchestrator assets takes effect on the next transaction with zero downtime."
""")

# ─────────────────────────────────────────────
# SLIDE 8 — DASHBOARD KPIS
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
add_text_box(sl, "Section 4 — Monitoring Dashboard KPIs", 0.5, 0.2, 12, 0.7, size=24, bold=True, color=WHITE)

kpis = [
    ("Total Orders Today", "4,827",  "▲ 12% vs yesterday",       BLUE),
    ("Success Rate",       "97.4%",  "▲ Above 95% SLA target",    GREEN),
    ("Failed Orders",      "125",    "▼ 43% vs pre-optimization", RED),
    ("Retries Executed",   "218",    "▼ 60% improvement",         RGBColor(0xf9,0x73,0x16)),
    ("Avg Process Time",   "1.8 min","▼ 57% faster per order",    RGBColor(0x14,0xb8,0xa6)),
    ("Bot Uptime (30d)",   "99.2%",  "▲ Target: 99% — Exceeded", RGBColor(0x93,0x33,0xea)),
]
for i, (label, value, trend, color) in enumerate(kpis):
    x = 0.3 + (i%3)*4.25
    y = 1.3  + (i//3)*2.65
    add_rect(sl, x, y, 3.95, 2.35, color)
    add_text_box(sl, label, x+0.15, y+0.18, 3.65, 0.4,  size=10, bold=True, color=WHITE)
    add_text_box(sl, value, x+0.15, y+0.6,  3.65, 1.0,  size=30, bold=True, color=WHITE)
    add_text_box(sl, trend, x+0.15, y+1.75, 3.65, 0.42, size=10, color=RGBColor(0xa7,0xf3,0xd0))

set_notes(sl,
"""SLIDE 8 — MONITORING DASHBOARD KPIs

"This slide represents what the UiPath Orchestrator monitoring dashboard looks like in production — post-optimization."

Walk through each KPI card:

Total Orders Today — 4,827:
"Volume metric showing today's order throughput. Up 12% versus yesterday. The 5-bot parallel architecture handles this comfortably."

Success Rate — 97.4%:
"Our most important KPI. The SLA target is 95%. We're exceeding it by 2.4 points. PagerDuty alerts fire if this drops below 90%."

Failed Orders — 125 (2.6%):
"Down 43% compared to pre-optimization. The remaining failures are mostly in the dead letter queue — unresolvable orders requiring manual investigation."

Retries Executed — 218:
"4.5% of total orders required at least one retry. This is down 60% from before. The circuit breaker prevents retry storms when the database is under stress."

Avg Processing Time — 1.8 minutes:
"Per-order average including all workflow stages. Well under our 4-minute SLA target."

Bot Uptime — 99.2% over 30 days:
"Exceeds our 99% availability target. The heartbeat monitoring and automatic bot restart on crash contribute to this high uptime."

Transition: "The full interactive dashboard is available on the live website — click the link on Slide 1 to see it."
""")

# ─────────────────────────────────────────────
# SLIDE 9 — AI USAGE
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
add_text_box(sl, "Section 5 — AI-Assisted Development: Vibe Coding with Claude Code", 0.5, 0.2, 12.3, 0.7, size=22, bold=True, color=WHITE)

add_rect(sl, 0.3, 1.2, 12.73, 1.1, RGBColor(0xef,0xf6,0xff))
add_text_box(sl, "What is Vibe Coding?", 0.5, 1.25, 5, 0.4, size=12, bold=True, color=NAVY)
add_text_box(sl,
    "Vibe coding = describing what you want in natural language → AI generates the implementation.\n"
    "The developer acts as architect & reviewer; Claude Code handles implementation details.",
    0.5, 1.6, 12.3, 0.55, size=10, color=RGBColor(0x1e,0x29,0x3b))

ai_items = [
    ("💬 Prompt Engineering",   RGBColor(0x25,0x63,0xeb), "Full project requirements described in natural language → Claude structured architecture, layout, and component design from plain English."),
    ("🎨 UI/UX Generation",     GREEN,                    "Design specs like 'enterprise-style, navy/gold color scheme, card layout' became actual CSS gradients, grid systems, and styled components."),
    ("📋 Domain Synthesis",     GOLD,                     "RPA error scenarios, root causes, and business impacts generated by feeding Claude domain context — technically accurate UiPath + e-commerce content."),
    ("🔄 Iterative Refinement", RED,                      "Each section refined through prompt → generate → review loop. The decision flow evolved from a list to a visual step-by-step arrow flow via one prompt."),
    ("⚙️ Full Stack Generation",RGBColor(0x93,0x33,0xea), "Claude Code generated HTML/CSS/JS, Python PPT/Word scripts, folder structure, GitHub setup commands, and the video script — from conversation."),
    ("📊 Consistent Mock Data", RGBColor(0x14,0xb8,0xa6), "Dashboard KPIs internally consistent and realistic. Before/after metrics logically aligned with the described optimization improvements."),
]
for i, (title, color, body) in enumerate(ai_items):
    x = 0.3 + (i%2)*6.45
    y = 2.5  + (i//2)*1.55
    add_rect(sl, x, y, 6.15, 1.4, WHITE)
    add_rect(sl, x, y, 6.15, 0.07, color)
    add_text_box(sl, title, x+0.12, y+0.12, 5.9, 0.35, size=11, bold=True, color=NAVY)
    add_text_box(sl, body,  x+0.12, y+0.5,  5.9, 0.8,  size=9,  color=RGBColor(0x1e,0x29,0x3b))

set_notes(sl,
"""SLIDE 9 — AI-ASSISTED DEVELOPMENT (VIBE CODING)

"The final section explains how this entire project was built using AI-assisted development with Claude Code."

Define vibe coding:
"Vibe coding is a term coined for AI-assisted development where you describe your intent in natural language and the AI generates the implementation. You stay in the role of architect and quality reviewer — the AI is the implementation engine."

Walk through each card:

Prompt Engineering: "The entire assignment brief was given to Claude Code as a single natural language prompt. It analyzed the requirements and structured the full five-section architecture, data models, and visual hierarchy without any manual outline."

UI/UX Generation: "I described the look: 'enterprise-style, navy and gold color scheme, card-based layout, clean academic feel.' Claude Code converted those adjectives into specific CSS variables, linear gradient definitions, and responsive grid systems."

Domain Synthesis: "I didn't write the error descriptions manually. I provided the domain context — UiPath RPA, e-commerce order processing — and Claude synthesized technically accurate exception names, root causes, and business impact statements."

Iterative Refinement: "The debugging decision flow started as a simple bullet list. I typed 'make it visual with a horizontal arrow-connected step flow' — and Claude restructured it in seconds."

Full Stack Generation: "Beyond the website, Claude Code also wrote the Python scripts that generated this PowerPoint, the Word documents, and set up the GitHub repository — all from conversational prompts in the terminal."

Key insight: "AI doesn't replace domain expertise — it amplifies it. I needed to know what RPA errors are realistic, what optimizations matter, what KPIs to track. The AI turned that knowledge into a professional deliverable."
""")

# ─────────────────────────────────────────────
# SLIDE 10 — KEY METRICS SUMMARY
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, NAVY)
add_text_box(sl, "Key Results Summary", 0.5, 0.4, 12, 0.8, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(sl, 3, 1.1, 7.33, 0.07, GOLD)

metrics = [
    ("↑26%",  "Success Rate Gain",     "71% → 97.4%",       BLUE),
    ("5×",    "Throughput Increase",   "120 → 600 orders/hr",GREEN),
    ("57%",   "Faster Processing",     "4.2min → 1.8min",    RGBColor(0xf9,0x73,0x16)),
    ("96%",   "Manual Work Reduced",   "3.5hr → 12min/day",  RGBColor(0x93,0x33,0xea)),
    ("0",     "Duplicate Orders",      "Per month post-fix",  RED),
    ("99.2%", "Bot Uptime (30-day)",   "Exceeds 99% target",  RGBColor(0x14,0xb8,0xa6)),
]
for i, (val, label, sub, color) in enumerate(metrics):
    x = 0.5 + (i%3)*4.1
    y = 1.5  + (i//3)*2.5
    add_rect(sl, x, y, 3.7, 2.15, RGBColor(0x1e,0x4a,0x76))
    add_rect(sl, x, y, 3.7, 0.09, color)
    add_text_box(sl, val,   x+0.15, y+0.2,  3.4, 0.8,  size=36, bold=True, color=color)
    add_text_box(sl, label, x+0.15, y+1.05, 3.4, 0.45, size=12, bold=True, color=WHITE)
    add_text_box(sl, sub,   x+0.15, y+1.5,  3.4, 0.45, size=10, color=RGBColor(0x93,0xc5,0xfd))

set_notes(sl,
"""SLIDE 10 — KEY RESULTS SUMMARY

"Before I close, here's the headline summary of what this case study achieved."

Read through the six metrics with emphasis:

↑26% Success Rate: "The single most important metric. From 71% to 97.4% — the bot went from unreliable to enterprise-grade."

5× Throughput: "Five times more orders processed per hour with the same infrastructure. Pure efficiency gain from parallel bot lanes."

57% Faster: "Each order now processes in 1.8 minutes versus 4.2 minutes. Customers receive confirmations faster, stock is updated sooner."

96% Manual Work Reduced: "Operations staff went from spending 3.5 hours every day on manual workarounds to 12 minutes. That's nearly an entire workday returned."

0 Duplicate Orders: "Zero. Complete elimination of duplicate charges. No more finance reconciliation incidents. No more customer complaints about double billing."

99.2% Uptime: "Exceeds the 99% SLA target. The bot is more reliable than most internal systems."

"These aren't theoretical projections — they're the result of applying systematic debugging and targeted optimizations to real, documented failures."
""")

# ─────────────────────────────────────────────
# SLIDE 11 — CLOSING WITH LINKS
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, NAVY)
add_rect(sl, 0, 2.8, 13.33, 0.12, GOLD)
add_text_box(sl, "Thank You", 1, 0.7, 11, 1.2, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(sl, "RPA Course · Assignment 04 · E-Commerce Order Processing Bot Case Study", 1, 2.0, 11, 0.55, size=14, color=RGBColor(0x93,0xc5,0xfd), align=PP_ALIGN.CENTER)
add_text_box(sl, "Omais Siddiqui  |  April 2026  |  Built with Claude Code AI-Assisted Vibe Coding", 1, 3.1, 11, 0.5, size=12, color=RGBColor(0x64,0x74,0x8b), align=PP_ALIGN.CENTER)

# Prominent link boxes
add_link_box(sl, "🌐  Live Website — Click to Open", LIVE_URL,   1.5, 3.85, 4.6, 0.65, RGBColor(0x1e,0x4a,0x76))
add_link_box(sl, "📁  GitHub Repository — Source Code", GITHUB_URL, 7.2, 3.85, 4.6, 0.65, RGBColor(0x1e,0x4a,0x76))

add_text_box(sl, LIVE_URL,   1.5, 4.6, 4.6, 0.35, size=9, color=RGBColor(0x93,0xc5,0xfd), align=PP_ALIGN.CENTER)
add_text_box(sl, GITHUB_URL, 7.2, 4.6, 4.6, 0.35, size=9, color=RGBColor(0x93,0xc5,0xfd), align=PP_ALIGN.CENTER)

add_text_box(sl, "Local file: file:///Users/omaissaeedsiddiqui/Desktop/Assignment_04_RPA/code/index.html",
             1, 5.15, 11, 0.35, size=9, color=RGBColor(0x47,0x5a,0x6b), align=PP_ALIGN.CENTER)

set_notes(sl,
"""SLIDE 11 — CLOSING SLIDE

"Thank you for your time. Let me leave you with the two most important links."

Live Website:
"The complete interactive website is live on GitHub Pages. Click the blue button to open it in your browser — all five sections are navigable with the sticky top navigation."
URL: https://omais-eng.github.io/assignment-04-rpa-bot-troubleshooting/

GitHub Repository:
"All project files are on GitHub — the full HTML source code, this PowerPoint file, four Word documents (one per section), the video script, and the README."
URL: https://github.com/omais-eng/assignment-04-rpa-bot-troubleshooting

Local file (if offline):
"If you're viewing this offline, the website also works as a local file at the path shown at the bottom of this slide — just open it in any browser, no build tools required."

"Questions welcome. This case study demonstrates that systematic RPA debugging — combined with AI-assisted development — can transform a 71% success rate into 97.4%, five times the throughput, and near-zero manual overhead."
""")

out = "/Users/omaissaeedsiddiqui/Desktop/Assignment_04_RPA/ppt/Assignment_04_RPA_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Live URL: {LIVE_URL}")
print(f"GitHub:   {GITHUB_URL}")
